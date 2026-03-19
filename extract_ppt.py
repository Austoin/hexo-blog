#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""提取PPT文本内容"""

from pptx import Presentation
import os
import sys


def extract_text_from_pptx(pptx_path):
    """从PPTX文件中提取所有文本"""
    try:
        prs = Presentation(pptx_path)
        text_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n{'=' * 60}\n第 {slide_num} 页\n{'=' * 60}\n"

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"

            if slide_text.strip():
                text_content.append(slide_text)

        return "\n".join(text_content)
    except Exception as e:
        return f"错误: {str(e)}"


def main():
    # PPT文件列表
    ppt_dir = r"C:\Users\liu'zhi'gui\Desktop\02 27考研王道计算机【数据结构领学班】\pdf"

    ppt_files = [
        "[11]--2.3.1_单链表的定义_20260316202516.pptx",
        "[12]--2.3.2_1_单链表的插入删除_20260316205919.pptx",
        "[13]--2.3.2_2_单链表的查找_20260316210547.pptx",
        "[14]--2.3.2_3_单链表的建立_20260316211338.pptx",
        "[15]--2.3.3_双链表_20260316211511.pptx",
        "[16]--2.3.4_循环链表_20260316213415.pptx",
        "[17]--2.3.5_静态链表_20260316233100.pptx",
        "[18]--2.3.6_顺序表和链表的比较_20260316234040.pptx",
    ]

    for ppt_file in ppt_files:
        ppt_path = os.path.join(ppt_dir, ppt_file)
        print(f"\n\n{'#' * 80}")
        print(f"文件: {ppt_file}")
        print(f"{'#' * 80}")

        if os.path.exists(ppt_path):
            content = extract_text_from_pptx(ppt_path)
            print(content)
        else:
            print(f"文件不存在: {ppt_path}")


if __name__ == "__main__":
    main()
